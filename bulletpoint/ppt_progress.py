# -*- coding: utf-8 -*-
import re
import sys
import nltk
import sys  
import random
from pptx import Presentation

bullet = '\u2022'
bullet = u'•'

def bulleted(string):
    return "%s %s" %(bullet, string)

def print_hierarchy(sentence_sets):
    for i in range(len(sentence_sets)):
        sentence_set = sentence_sets[i]
        for j in range(0, len(sentence_set)):
            current_sentence = sentence_set[j]
            if current_sentence != "":
                # Reconstruct sentence to encode hierarchy
                that_indexes = [m.start() for m in re.finditer('that', current_sentence)]
                new_sentence = []
                for i in range(len(current_sentence)):
                    if i in that_indexes:
                        new_sentence.append("\n\t\t %s %s" %(bullet, current_sentence[i]))
                    else:
                        new_sentence.append(current_sentence[i])
                new_sentence = ''.join(new_sentence)
                if j == 0:
                    print new_sentence
                else:   
                    print "\t%s" %(bulleted(new_sentence)) 
        print ""

def add_bulleted_slides(sentence_sets, presentation):
    for i in range(len(sentence_sets)):
        sentence_set = sentence_sets[i]
        words = ' '.join(sentence_set)
        tokenized_words = nltk.tokenize.word_tokenize(words)

        bullet_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        # Determine slide title: the trigram with highest PMI in paragraph
        trigram_measures = nltk.collocations.TrigramAssocMeasures()
        finder = nltk.collocations.TrigramCollocationFinder.from_words(tokenized_words)
        best_ngram = list(finder.nbest(trigram_measures.pmi, 1)[0])
        title_shape.text = ' '.join(best_ngram)
        
        tf = body_shape.text_frame
        tf.text = sentence_set[0]

        for j in range(1, len(sentence_set)):
            current_sentence = sentence_set[j]
            # Reconstruct sentence to encode hierarchy
            that_indexes = [m.start() for m in re.finditer('that', current_sentence)]
            for k in range(len(that_indexes)-1):
                that_index = that_indexes[k]
                substring = sentence[that_index:that_index[k+1]+1]
                p = tf.add_paragraph()
                p.text = substring
                p.level = 1

    return presentation

# I want a function that given a string, will return a sequential list of tuples 
# that each consist of the sentence and the level
def get_hierarchy(sentence, level):
    bulletpoint_list = []
    that_index = sentence.find('that')
    if that_index == -1:
        return [(sentence, level)]
    else:
        new_bullet = sentence[:that_index]
        bulletpoint_list.append((new_bullet, level))
        level += int(round(random.random()))
        next_section = sentence[that_index+4:]
        return bulletpoint_list + get_hierarchy(next_section, level)

if __name__ == "__main__":
    sentence = "I got to go to school that is in San Diego and that is the best ever and which I love which I love which I love which I love."
    level = 1
    print(get_hierarchy(sentence, level))

    # if len(sys.argv) < 1:
    #     print("You must call program as: python ppt.py <datafile> <'prose'|'poetry'>")
    #     print("This program generates a powerpoint file from the text in datafile.")
    #     sys.exit(1)
    # filename = sys.argv[1]
    # content_type = sys.argv[2]

    # filename = 'mit_news_3.txt'


    # f = open(filename, 'r')
    # content = f.read().decode('utf-8').strip()
    # paragraphs = content.split("\n")
    # tokenized_sentence_sets = [nltk.tokenize.sent_tokenize(paragraph) for paragraph in paragraphs]
    # print_hierarchy(tokenized_sentence_sets)

    # presentation = Presentation()
    # title_slide_layout = presentation.slide_layouts[0]
    # slide = presentation.slides.add_slide(title_slide_layout)
    # title = slide.shapes.title
    # subtitle = slide.placeholders[1]

    # title.text = filename
    # subtitle.text = "python-pptx was here!" # trigram with highest PMI in entire source
    # presentation = add_bulleted_slides(tokenized_sentence_sets, presentation)
    # presentation.save('test.pptx')

