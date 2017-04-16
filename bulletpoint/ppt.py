# -*- coding: utf-8 -*-
import re
import sys
import nltk
import random
from pptx import Presentation


def get_hierarchy(sentence, level):
    bulletpoint_list = []
    that_index = sentence.find('that')
    if that_index == -1:
        return [(sentence, level)]
    else:
        new_bullet = sentence[:that_index]
        bulletpoint_list.append((new_bullet, level))
        level += int(round(random.random()))
        next_section = sentence[that_index+4:]
        return bulletpoint_list + get_hierarchy(next_section, level) 

def get_best_trigrams(text, num_best):
        tokenized_words = nltk.tokenize.word_tokenize(text)
        trigram_measures = nltk.collocations.TrigramAssocMeasures()
        finder = nltk.collocations.TrigramCollocationFinder.from_words(tokenized_words)
        best_ngrams = finder.nbest(trigram_measures.pmi, num_best)
        to_return =  [' '.join(list(best_ngram)) for best_ngram in best_ngrams]
        return to_return

def add_bulleted_slides(sentence_sets, presentation):
    for i in range(len(sentence_sets)):
        sentence_set = sentence_sets[i]
        words = ' '.join(sentence_set)
        

        bullet_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        # Determine slide title: the trigram with highest PMI in paragraph
        title_shape.text = get_best_trigrams(words, 1)[0]
        
        tf = body_shape.text_frame
        tf.text = sentence_set[0]

        for j in range(1, len(sentence_set)):
            current_sentence = sentence_set[j]
            # Reconstruct sentence to encode hierarchy
            level = 1
            phrases_with_levels = get_hierarchy(current_sentence, level)
            for phrase, level in phrases_with_levels:
                p = tf.add_paragraph()
                p.text = phrase
                p.level = level
    return presentation

def get_chunks(content, opt_chunk_size=None, poem=False):
    if opt_chunk_size:
        if poem:
            lines = content.split("\n")
            content = '. '.join(lines)
        sentences = nltk.tokenize.sent_tokenize(content)
        tokenized_sentence_sets = [sentences[i:i + opt_chunk_size] for i in xrange(0, len(sentences), opt_chunk_size)]
    else:
        paragraphs = content.split("\n")
        tokenized_sentence_sets = [nltk.tokenize.sent_tokenize(paragraph) for paragraph in paragraphs]
    return tokenized_sentence_sets

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("This program generates a powerpoint file from the text in datafile.")
        print("You must call program as: python ppt.py <path_to_file> <content_type>")
        print("content_type can be prose or poem")
        sys.exit(1)
    filename = sys.argv[1]
    content_type = sys.argv[2]
    poem = content_type == "poem"
    if poem:
        chunk_size = 4
    else:
        chunk_size = 2

    f = open(filename, 'r')
    content = f.read().decode('utf-8').strip()
    tokenized_sentence_sets = get_chunks(content, chunk_size, poem)

    presentation = Presentation()
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    best = get_best_trigrams(content, 2)
    title.text = best[0]
    subtitle.text = best[1] # trigram with highest PMI in entire source
    presentation = add_bulleted_slides(tokenized_sentence_sets, presentation)
    presentation.save(filename + '.pptx')

